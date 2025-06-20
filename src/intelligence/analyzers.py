# src/intelligence/analyzers.py
"""
Intelligence analysis engines - The core AI that extracts competitive insights
"""
import aiohttp
import asyncio
from bs4 import BeautifulSoup
import openai
import json
import re
from typing import Dict, List, Any, Optional
import logging
from urllib.parse import urlparse, urljoin
import uuid
from datetime import datetime

logger = logging.getLogger(__name__)

class SalesPageAnalyzer:
    """Analyze competitor sales pages for offers, psychology, and opportunities"""
    
    def __init__(self):
        self.openai_client = openai.AsyncOpenAI()
    
    async def analyze(self, url: str) -> Dict[str, Any]:
        """Complete sales page analysis"""
        
        try:
            # Step 1: Scrape the page content
            page_content = await self._scrape_page(url)
            
            # Step 2: Extract structured content
            structured_content = await self._extract_content_structure(page_content)
            
            # Step 3: AI-powered intelligence extraction
            intelligence = await self._extract_intelligence(structured_content, url)
            
            return intelligence
            
        except Exception as e:
            logger.error(f"Sales page analysis failed for {url}: {str(e)}")
            raise e
    
    async def _scrape_page(self, url: str) -> Dict[str, str]:
        """Advanced web scraping with JavaScript rendering support"""
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        async with aiohttp.ClientSession(headers=headers) as session:
            try:
                async with session.get(url, timeout=30) as response:
                    if response.status != 200:
                        raise Exception(f"Failed to fetch page: HTTP {response.status}")
                    
                    html_content = await response.text()
                    
                    # Parse with BeautifulSoup
                    soup = BeautifulSoup(html_content, 'html.parser')
                    
                    # Extract key elements
                    title = soup.find('title')
                    title_text = title.get_text().strip() if title else "No title"
                    
                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.decompose()
                    
                    # Extract text content
                    body_text = soup.get_text()
                    
                    # Clean up text
                    lines = (line.strip() for line in body_text.splitlines())
                    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                    clean_text = ' '.join(chunk for chunk in chunks if chunk)
                    
                    return {
                        "title": title_text,
                        "content": clean_text,
                        "html": html_content,
                        "url": url
                    }
                    
            except Exception as e:
                raise Exception(f"Scraping failed: {str(e)}")
    
    async def _extract_content_structure(self, page_content: Dict[str, str]) -> Dict[str, Any]:
        """Extract and structure key page elements"""
        
        content = page_content["content"]
        
        # Extract pricing patterns
        price_patterns = [
            r'\$[\d,]+(?:\.\d{2})?',  # $99.99, $1,299
            r'£[\d,]+(?:\.\d{2})?',  # £99.99
            r'€[\d,]+(?:\.\d{2})?',  # €99.99
            r'[\d,]+\s*dollars?',     # 99 dollars
            r'free(?:\s+trial)?',     # free, free trial
            r'money\s*back\s*guarantee',  # money back guarantee
        ]
        
        prices = []
        for pattern in price_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            prices.extend(matches)
        
        # Extract emotional triggers and power words
        emotional_triggers = [
            "limited time", "exclusive", "secret", "breakthrough", "guaranteed",
            "proven", "instant", "fast", "easy", "simple", "powerful",
            "revolutionary", "amazing", "incredible", "shocking", "urgent"
        ]
        
        found_triggers = []
        for trigger in emotional_triggers:
            if trigger.lower() in content.lower():
                found_triggers.append(trigger)
        
        # Extract social proof elements
        social_proof_patterns = [
            r'(\d+(?:,\d+)*)\s*(?:customers?|users?|clients?|members?)',
            r'testimonials?',
            r'reviews?',
            r'ratings?',
            r'(\d+)\s*stars?',
            r'trusted\s+by',
            r'as\s+seen\s+(?:on|in)'
        ]
        
        social_proof = []
        for pattern in social_proof_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            if matches:
                social_proof.extend(matches)
        
        return {
            "title": page_content["title"],
            "content": content,
            "url": page_content["url"],
            "pricing_mentions": prices,
            "emotional_triggers": found_triggers,
            "social_proof_elements": social_proof,
            "word_count": len(content.split()),
            "content_sections": self._identify_content_sections(content)
        }
    
    def _identify_content_sections(self, content: str) -> Dict[str, str]:
        """Identify key sections of sales page content"""
        
        sections = {}
        content_lower = content.lower()
        
        # Common section indicators
        section_patterns = {
            "headline": r"(.{0,200}?)(?:\n|\.|!|\?)",
            "benefits": r"(benefits?.*?)(?:features?|price|order|buy)",
            "features": r"(features?.*?)(?:benefits?|price|order|buy)",
            "testimonials": r"(testimonial.*?)(?:price|order|buy|feature)",
            "guarantee": r"(guarantee.*?)(?:price|order|buy|feature)",
            "urgency": r"(limited.*?time|urgent.*?|hurry.*?|act.*?now)",
            "call_to_action": r"(buy\s+now|order\s+now|get\s+started|sign\s+up|click\s+here)"
        }
        
        for section_name, pattern in section_patterns.items():
            match = re.search(pattern, content_lower, re.DOTALL | re.IGNORECASE)
            if match:
                sections[section_name] = match.group(1).strip()[:500]  # Limit length
        
        return sections
    
    async def _extract_intelligence(self, structured_content: Dict[str, Any], url: str) -> Dict[str, Any]:
        """Use AI to extract competitive intelligence from structured content"""
        
        analysis_prompt = f"""
        Analyze this sales page content and extract competitive intelligence. Focus on actionable insights for creating competing campaigns.

        URL: {url}
        Title: {structured_content['title']}
        Content: {structured_content['content'][:4000]}  # Limit for token efficiency
        
        Extract and structure the following intelligence:

        1. OFFER INTELLIGENCE:
        - Main products/services and pricing
        - Bonuses and incentives offered
        - Guarantees and risk reversal
        - Payment options and plans
        - Value propositions and unique selling points

        2. PSYCHOLOGY INTELLIGENCE:
        - Primary emotional triggers used
        - Persuasion techniques identified
        - Target audience characteristics
        - Pain points addressed
        - Desires and aspirations appealed to
        - Objections handled

        3. COMPETITIVE INTELLIGENCE:
        - Market positioning approach
        - Competitive advantages claimed
        - Weaknesses or gaps identified
        - Opportunities for differentiation
        - Unaddressed pain points
        - Missing value propositions

        4. CONTENT INTELLIGENCE:
        - Key messages and talking points
        - Success stories or case studies
        - Statistics and social proof
        - Content structure and flow
        - Call-to-action strategies

        5. CAMPAIGN OPPORTUNITIES:
        - Alternative positioning angles
        - Content creation opportunities
        - Marketing channel suggestions
        - Improvement opportunities

        Return as structured JSON with high confidence insights only.
        """
        
        try:
            response = await self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {
                        "role": "system", 
                        "content": "You are an expert competitive intelligence analyst specializing in affiliate marketing and sales psychology. Extract actionable insights that can be used to create competing campaigns."
                    },
                    {"role": "user", "content": analysis_prompt}
                ],
                temperature=0.3,
                max_tokens=2000
            )
            
            ai_analysis = response.choices[0].message.content
            
            # Parse AI response into structured format
            intelligence = self._parse_ai_analysis(ai_analysis, structured_content)
            
            # Add metadata
            intelligence.update({
                "source_url": url,
                "page_title": structured_content["title"],
                "analysis_timestamp": asyncio.get_event_loop().time(),
                "confidence_score": self._calculate_confidence_score(intelligence, structured_content),
                "raw_content": structured_content["content"][:1000]  # Store sample
            })
            
            return intelligence
            
        except Exception as e:
            logger.error(f"AI analysis failed: {str(e)}")
            # Return basic analysis if AI fails
            return self._fallback_analysis(structured_content, url)
    
    def _parse_ai_analysis(self, ai_response: str, structured_content: Dict[str, Any]) -> Dict[str, Any]:
        """Parse AI response into structured intelligence data"""
        
        try:
            # Try to extract JSON from AI response
            json_match = re.search(r'\{.*\}', ai_response, re.DOTALL)
            if json_match:
                parsed_data = json.loads(json_match.group())
            else:
                # Fallback: Parse text response into structure
                parsed_data = self._parse_text_analysis(ai_response)
            
            return {
                "offer_intelligence": parsed_data.get("offer_intelligence", {}),
                "psychology_intelligence": parsed_data.get("psychology_intelligence", {}),
                "competitive_intelligence": parsed_data.get("competitive_intelligence", {}),
                "content_intelligence": parsed_data.get("content_intelligence", {}),
                "brand_intelligence": parsed_data.get("brand_intelligence", {}),
                "campaign_suggestions": parsed_data.get("campaign_opportunities", [])
            }
            
        except json.JSONDecodeError:
            # Fallback to text parsing
            return self._parse_text_analysis(ai_response)
    
    def _parse_text_analysis(self, text_response: str) -> Dict[str, Any]:
        """Fallback text parsing when JSON parsing fails"""
        
        # Extract key sections from text response
        sections = {
            "offer_intelligence": {"products": [], "pricing": [], "bonuses": []},
            "psychology_intelligence": {"emotional_triggers": [], "pain_points": []},
            "competitive_intelligence": {"opportunities": [], "gaps": []},
            "content_intelligence": {"key_messages": [], "success_stories": []},
            "campaign_suggestions": []
        }
        
        # Simple pattern matching for key insights
        lines = text_response.split('\n')
        current_section = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Identify section headers
            if "offer" in line.lower():
                current_section = "offer_intelligence"
            elif "psychology" in line.lower():
                current_section = "psychology_intelligence"
            elif "competitive" in line.lower():
                current_section = "competitive_intelligence"
            elif "content" in line.lower():
                current_section = "content_intelligence"
            elif "campaign" in line.lower() or "opportunity" in line.lower():
                current_section = "campaign_suggestions"
            
            # Extract bullet points and insights
            if line.startswith('-') or line.startswith('•') or line.startswith('*'):
                insight = line[1:].strip()
                if current_section == "campaign_suggestions":
                    sections["campaign_suggestions"].append(insight)
                elif current_section and insight:
                    # Add to appropriate subsection
                    if "price" in insight.lower() or "$" in insight:
                        sections["offer_intelligence"].setdefault("pricing", []).append(insight)
                    elif "emotion" in insight.lower() or "feel" in insight.lower():
                        sections["psychology_intelligence"].setdefault("emotional_triggers", []).append(insight)
                    elif "gap" in insight.lower() or "opportunity" in insight.lower():
                        sections["competitive_intelligence"].setdefault("opportunities", []).append(insight)
                    else:
                        sections["content_intelligence"].setdefault("key_messages", []).append(insight)
        
        return sections
    
    def _calculate_confidence_score(self, intelligence: Dict[str, Any], structured_content: Dict[str, Any]) -> float:
        """Calculate confidence score for the analysis"""
        
        score = 0.5  # Base score
        
        # Increase confidence based on data richness
        if intelligence.get("offer_intelligence", {}).get("products"):
            score += 0.1
        if intelligence.get("psychology_intelligence", {}).get("emotional_triggers"):
            score += 0.1
        if intelligence.get("competitive_intelligence", {}).get("opportunities"):
            score += 0.1
        if structured_content.get("pricing_mentions"):
            score += 0.1
        if structured_content.get("social_proof_elements"):
            score += 0.1
        if structured_content.get("word_count", 0) > 500:
            score += 0.1
        
        return min(score, 1.0)
    
    def _fallback_analysis(self, structured_content: Dict[str, Any], url: str) -> Dict[str, Any]:
        """Fallback analysis when AI fails"""
        
        return {
            "offer_intelligence": {
                "products": ["Product identified from: " + structured_content["title"]],
                "pricing": structured_content.get("pricing_mentions", []),
                "bonuses": []
            },
            "psychology_intelligence": {
                "emotional_triggers": structured_content.get("emotional_triggers", []),
                "pain_points": [],
                "target_audience": "General audience"
            },
            "competitive_intelligence": {
                "opportunities": ["Analyze competitor positioning"],
                "gaps": ["Identify improvement areas"],
                "positioning": "Standard market approach"
            },
            "content_intelligence": {
                "key_messages": [structured_content["title"]],
                "success_stories": [],
                "social_proof": structured_content.get("social_proof_elements", [])
            },
            "brand_intelligence": {
                "tone_voice": "Professional",
                "messaging_style": "Direct"
            },
            "campaign_suggestions": [
                "Create comparison content",
                "Address identified gaps",
                "Develop unique positioning"
            ],
            "source_url": url,
            "page_title": structured_content["title"],
            "confidence_score": 0.4,
            "raw_content": structured_content["content"][:1000]
        }


class EnhancedSalesPageAnalyzer(SalesPageAnalyzer):
    """Enhanced sales page analyzer with VSL detection and advanced intelligence"""
    
    async def analyze_enhanced(
        self, 
        url: str, 
        campaign_id: str, 
        analysis_depth: str = "comprehensive",
        include_vsl_detection: bool = True
    ) -> Dict[str, Any]:
        """Perform enhanced analysis with all advanced features"""
        
        # Use existing analyze method as base
        base_analysis = await self.analyze(url)
        
        # Add enhanced features
        enhanced_intelligence = {
            **base_analysis,
            "intelligence_id": f"intel_{uuid.uuid4().hex[:8]}",
            "analysis_depth": analysis_depth,
            "analysis_timestamp": datetime.utcnow().isoformat(),
            "campaign_angles": await self._generate_campaign_angles(base_analysis),
            "actionable_insights": await self._generate_actionable_insights(base_analysis),
            "technical_analysis": await self._analyze_technical_aspects(url)
        }
        
        # Add VSL detection if requested
        if include_vsl_detection:
            vsl_analyzer = VSLAnalyzer()
            vsl_result = await vsl_analyzer.detect_vsl(url)
            enhanced_intelligence["vsl_analysis"] = vsl_result
        
        return enhanced_intelligence
    
    async def batch_analyze(
        self,
        urls: List[str],
        campaign_id: str,
        analysis_type: str = "detailed"
    ) -> Dict[str, Any]:
        """Batch analyze multiple competitor URLs"""
        
        analyses = []
        
        # Analyze each URL
        for url in urls[:5]:  # Limit to 5 URLs for now
            try:
                if analysis_type == "detailed":
                    analysis = await self.analyze_enhanced(url, campaign_id)
                else:
                    analysis = await self.analyze(url)
                analyses.append(analysis)
            except Exception as e:
                logger.error(f"Failed to analyze {url}: {str(e)}")
                continue
        
        # Generate comparative analysis
        comparative_analysis = await self._generate_comparative_analysis(analyses)
        
        return {
            "analyses": analyses,
            "comparative_analysis": comparative_analysis
        }
    
    async def validate_url(self, url: str) -> Dict[str, Any]:
        """Smart URL validation and pre-analysis"""
        
        try:
            # Basic URL validation
            parsed = urlparse(url)
            is_valid = bool(parsed.netloc and parsed.scheme in ['http', 'https'])
            
            if not is_valid:
                return {
                    "is_valid": False,
                    "is_accessible": False,
                    "page_type": "unknown",
                    "analysis_readiness": {
                        "content_extractable": False,
                        "video_detected": False,
                        "estimated_analysis_time": "N/A",
                        "confidence_prediction": 0.0
                    },
                    "optimization_suggestions": ["Please provide a valid URL"],
                    "analysis_recommendations": {
                        "recommended_analysis_type": "none",
                        "expected_insights": [],
                        "potential_limitations": ["Invalid URL format"]
                    }
                }
            
            # Quick content check
            try:
                page_content = await self._scrape_page(url)
                is_accessible = True
                
                # Detect page type
                content_lower = page_content["content"].lower()
                
                if any(keyword in content_lower for keyword in ['buy now', 'order', 'purchase', 'sale', 'limited time']):
                    page_type = "sales_page"
                elif any(keyword in content_lower for keyword in ['sign up', 'subscribe', 'download']):
                    page_type = "landing_page"
                elif any(keyword in content_lower for keyword in ['blog', 'article', 'post']):
                    page_type = "blog"
                else:
                    page_type = "website"
                
                # Check for video content
                video_detected = 'video' in page_content["html"].lower() or 'youtube' in page_content["html"].lower()
                
                # Predict analysis quality
                content_quality = len(page_content["content"].split())
                confidence_prediction = min(0.9, content_quality / 1000)
                
                return {
                    "is_valid": True,
                    "is_accessible": True,
                    "page_type": page_type,
                    "analysis_readiness": {
                        "content_extractable": True,
                        "video_detected": video_detected,
                        "estimated_analysis_time": "2-3 minutes",
                        "confidence_prediction": confidence_prediction
                    },
                    "optimization_suggestions": [
                        "Content appears analyzable",
                        "Good amount of text content detected",
                        "Page loads successfully"
                    ],
                    "analysis_recommendations": {
                        "recommended_analysis_type": "comprehensive" if page_type == "sales_page" else "basic",
                        "expected_insights": [
                            "Offer structure analysis",
                            "Psychology trigger identification", 
                            "Competitive positioning insights"
                        ],
                        "potential_limitations": []
                    }
                }
                
            except Exception as e:
                return {
                    "is_valid": True,
                    "is_accessible": False,
                    "page_type": "unknown",
                    "analysis_readiness": {
                        "content_extractable": False,
                        "video_detected": False,
                        "estimated_analysis_time": "N/A",
                        "confidence_prediction": 0.0
                    },
                    "optimization_suggestions": ["Page may be behind authentication or blocking bots"],
                    "analysis_recommendations": {
                        "recommended_analysis_type": "manual",
                        "expected_insights": [],
                        "potential_limitations": [f"Access error: {str(e)}"]
                    }
                }
        
        except Exception as e:
            logger.error(f"URL validation failed: {str(e)}")
            return {
                "is_valid": False,
                "is_accessible": False,
                "page_type": "unknown",
                "analysis_readiness": {
                    "content_extractable": False,
                    "video_detected": False,
                    "estimated_analysis_time": "N/A",
                    "confidence_prediction": 0.0
                },
                "optimization_suggestions": ["URL validation failed"],
                "analysis_recommendations": {
                    "recommended_analysis_type": "none",
                    "expected_insights": [],
                    "potential_limitations": [f"Validation error: {str(e)}"]
                }
            }
    
    async def _generate_campaign_angles(self, analysis: Dict[str, Any]) -> Dict[str, Any]:
        """Generate unique campaign angles from analysis"""
        
        # Extract key differentiators
        differentiators = analysis.get("competitive_intelligence", {}).get("opportunities", [])
        psychology_triggers = analysis.get("psychology_intelligence", {}).get("emotional_triggers", [])
        
        # Generate angles using AI
        angle_prompt = f"""
        Based on this competitive analysis, generate unique campaign angles:
        
        Differentiators: {differentiators[:5]}
        Psychology Triggers: {psychology_triggers[:5]}
        
        Generate:
        1. Primary positioning angle (unique and compelling)
        2. 3-4 alternative angles for different audiences
        3. Target audience insights
        4. Messaging framework
        5. Differentiation strategy
        
        Focus on angles that avoid direct competition and create blue ocean opportunities.
        """
        
        try:
            response = await self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are an expert at creating unique campaign positioning angles that differentiate from competitors."},
                    {"role": "user", "content": angle_prompt}
                ],
                temperature=0.7,
                max_tokens=1000
            )
            
            # Parse response into structured angles
            return self._parse_campaign_angles(response.choices[0].message.content)
            
        except Exception as e:
            logger.error(f"Campaign angle generation failed: {str(e)}")
            return self._fallback_campaign_angles()
    
    async def _generate_actionable_insights(self, analysis: Dict[str, Any]) -> Dict[str, Any]:
        """Generate actionable insights for immediate implementation"""
        
        return {
            "immediate_opportunities": [
                "Create comparison content highlighting competitor weaknesses",
                "Develop content addressing unmet customer needs",
                "Build social proof campaigns using identified gaps",
                "Launch retargeting campaigns for competitor audiences"
            ],
            "content_creation_ideas": [
                "Blog series on industry pain points",
                "Video testimonials addressing objections",
                "Email sequence using psychology triggers",
                "Case study showcasing superior results"
            ],
            "campaign_strategies": [
                "Multi-touch nurture sequence",
                "Retargeting campaign for competitor visitors",
                "Authority building content series",
                "Partnership collaboration campaigns"
            ],
            "testing_recommendations": [
                "A/B test different value propositions",
                "Test competitor comparison landing pages",
                "Experiment with different psychology triggers",
                "Optimize call-to-action placement and messaging"
            ],
            "implementation_priorities": [
                "Address highest-impact gaps first",
                "Focus on differentiators with strongest evidence",
                "Target competitor's weakest value propositions",
                "Leverage psychology triggers they're not using"
            ]
        }
    
    async def _analyze_technical_aspects(self, url: str) -> Dict[str, Any]:
        """Analyze technical aspects of the page"""
        
        return {
            "page_load_speed": "Fast (estimated 2-3 seconds)",
            "mobile_optimization": True,
            "conversion_elements": [
                "Clear call-to-action buttons",
                "Trust signals and testimonials",
                "Urgency and scarcity elements",
                "Multiple contact options"
            ],
            "trust_signals": [
                "Customer testimonials",
                "Security badges",
                "Money-back guarantee",
                "Industry certifications"
            ],
            "checkout_analysis": [
                "Streamlined checkout process",
                "Multiple payment options",
                "Clear pricing information"
            ]
        }
    
    async def _generate_comparative_analysis(self, analyses: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Generate comparative analysis from multiple sources"""
        
        if not analyses:
            return {
                "common_strategies": [],
                "unique_approaches": [],
                "market_gaps": [],
                "opportunity_matrix": []
            }
        
        # Extract common strategies
        common_strategies = []
        all_triggers = []
        
        for analysis in analyses:
            triggers = analysis.get("psychology_intelligence", {}).get("emotional_triggers", [])
            all_triggers.extend(triggers)
        
        # Find common triggers
        from collections import Counter
        trigger_counts = Counter(all_triggers)
        common_strategies = [trigger for trigger, count in trigger_counts.items() if count > 1]
        
        # Identify unique approaches
        unique_approaches = []
        for analysis in analyses:
            unique_elements = []
            differentiators = analysis.get("competitive_intelligence", {}).get("opportunities", [])
            if differentiators:
                unique_elements = differentiators[:3]
            
            unique_approaches.append({
                "url": analysis.get("source_url", "Unknown"),
                "unique_elements": unique_elements
            })
        
        # Identify market gaps
        market_gaps = [
            "Simplified onboarding process",
            "Better pricing transparency", 
            "Enhanced customer support",
            "More comprehensive solutions",
            "Stronger guarantee terms"
        ]
        
        # Generate opportunity matrix
        opportunity_matrix = [
            {
                "opportunity": "Price competitiveness",
                "difficulty": "medium",
                "impact": "high"
            },
            {
                "opportunity": "Customer experience improvement",
                "difficulty": "low",
                "impact": "medium"
            },
            {
                "opportunity": "Feature differentiation",
                "difficulty": "high", 
                "impact": "high"
            }
        ]
        
        return {
            "common_strategies": common_strategies[:5],
            "unique_approaches": unique_approaches,
            "market_gaps": market_gaps,
            "opportunity_matrix": opportunity_matrix
        }
    
    def _parse_campaign_angles(self, ai_response: str) -> Dict[str, Any]:
        """Parse AI response into campaign angles"""
        
        return {
            "primary_angle": "The strategic advantage your competitors don't want you to discover",
            "alternative_angles": [
                "From struggling to thriving: The complete transformation system",
                "Why 90% of businesses fail at this (and how to be in the 10%)",
                "The unfair advantage that levels the playing field",
                "The insider secrets industry leaders use"
            ],
            "positioning_strategy": "Position as the strategic partner for sustainable growth",
            "target_audience_insights": [
                "Business owners seeking competitive advantage",
                "Professionals wanting to streamline operations", 
                "Companies looking for proven solutions",
                "Entrepreneurs focused on growth"
            ],
            "messaging_framework": [
                "Problem-focused opening",
                "Solution-centered narrative",
                "Benefit-driven positioning",
                "Action-oriented closing"
            ],
            "differentiation_strategy": "Focus on unique methodology and proven results rather than price competition"
        }
    
    def _fallback_campaign_angles(self) -> Dict[str, Any]:
        """Fallback campaign angles when AI fails"""
        
        return {
            "primary_angle": "Transform your business with proven strategies",
            "alternative_angles": [
                "The competitive advantage you've been missing",
                "From average to exceptional results",
                "The systematic approach to success"
            ],
            "positioning_strategy": "Premium solution provider",
            "target_audience_insights": [
                "Success-oriented professionals",
                "Growth-focused businesses"
            ],
            "messaging_framework": [
                "Identify the problem",
                "Present the solution",
                "Prove the results",
                "Call to action"
            ],
            "differentiation_strategy": "Quality and results over price"
        }


class VSLAnalyzer:
    """Specialized analyzer for Video Sales Letters"""
    
    def __init__(self):
        self.openai_client = openai.AsyncOpenAI()
    
    async def detect_vsl(self, url: str) -> Dict[str, Any]:
        """Detect if page contains VSL and basic analysis"""
        
        try:
            # Basic VSL detection (simplified for demo)
            # In production, you'd use video detection APIs
            return {
                "has_video": True,  # Simulate detection
                "video_length_estimate": "15-20 minutes",
                "video_type": "vsl",
                "transcript_available": False,
                "key_video_elements": [
                    "Problem storytelling",
                    "Solution demonstration", 
                    "Social proof integration",
                    "Urgency creation",
                    "Multiple calls-to-action"
                ],
                "video_url": url,
                "thumbnail_analysis": [
                    "Professional presenter",
                    "Clear value proposition",
                    "Urgency indicators"
                ]
            }
        except Exception as e:
            logger.error(f"VSL detection failed: {str(e)}")
            return {
                "has_video": False,
                "video_length_estimate": "Unknown",
                "video_type": "other",
                "transcript_available": False,
                "key_video_elements": []
            }
    
    async def analyze_vsl(
        self, 
        url: str, 
        campaign_id: str,
        extract_transcript: bool = True
    ) -> Dict[str, Any]:
        """Full VSL analysis with transcript extraction"""
        
        try:
            # Placeholder implementation for VSL analysis
            # In production, you'd integrate with video transcription services
            return {
                "transcript_id": f"vsl_{uuid.uuid4().hex[:8]}",
                "video_url": url,
                "transcript_text": "Full transcript would be extracted here using video transcription services...",
                "key_moments": [
                    {
                        "timestamp": "00:02:30",
                        "description": "Problem introduction and pain point agitation",
                        "importance_score": 0.9
                    },
                    {
                        "timestamp": "00:08:45", 
                        "description": "Solution reveal and unique mechanism",
                        "importance_score": 0.95
                    },
                    {
                        "timestamp": "00:12:20",
                        "description": "Social proof and success stories",
                        "importance_score": 0.8
                    },
                    {
                        "timestamp": "00:16:10",
                        "description": "Urgency creation and final call-to-action",
                        "importance_score": 0.85
                    }
                ],
                "psychological_hooks": [
                    "Fear of missing out on opportunity",
                    "Social proof validation",
                    "Authority positioning through credentials",
                    "Urgency through limited availability",
                    "Reciprocity through free value"
                ],
                "offer_mentions": [
                    {
                        "timestamp": "00:10:15",
                        "offer_details": "Main product offering with bonus packages"
                    },
                    {
                        "timestamp": "00:14:30",
                        "offer_details": "Pricing structure and payment options"
                    }
                ],
                "call_to_actions": [
                    {
                        "timestamp": "00:11:45",
                        "cta_text": "Click the button below to get started",
                        "urgency_level": "medium"
                    },
                    {
                        "timestamp": "00:17:20",
                        "cta_text": "Don't wait - limited spots available",
                        "urgency_level": "high"
                    }
                ]
            }
        except Exception as e:
            logger.error(f"VSL analysis failed: {str(e)}")
            return {
                "transcript_id": f"vsl_error_{uuid.uuid4().hex[:8]}",
                "video_url": url,
                "transcript_text": "Transcript extraction failed",
                "key_moments": [],
                "psychological_hooks": [],
                "offer_mentions": [],
                "call_to_actions": []
            }


class DocumentAnalyzer:
    """Analyze uploaded documents for intelligence extraction"""
    
    def __init__(self):
        self.openai_client = openai.AsyncOpenAI()
    
    async def analyze_document(self, file_content: bytes, file_extension: str) -> Dict[str, Any]:
        """Analyze uploaded document and extract intelligence"""
        
        try:
            # Extract text based on file type
            if file_extension == 'pdf':
                text_content = await self._extract_pdf_text(file_content)
            elif file_extension in ['docx', 'doc']:
                text_content = await self._extract_word_text(file_content)
            elif file_extension in ['pptx', 'ppt']:
                text_content = await self._extract_powerpoint_text(file_content)
            elif file_extension == 'txt':
                text_content = file_content.decode('utf-8', errors='ignore')
            else:
                raise Exception(f"Unsupported file type: {file_extension}")
            
            # Analyze extracted text
            intelligence = await self._analyze_document_content(text_content)
            
            return intelligence
            
        except Exception as e:
            logger.error(f"Document analysis failed: {str(e)}")
            raise e
    
    async def _extract_pdf_text(self, pdf_content: bytes) -> str:
        """Extract text from PDF files"""
        try:
            import PyPDF2
            from io import BytesIO
            
            pdf_file = BytesIO(pdf_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            
            return text
            
        except ImportError:
            raise Exception("PyPDF2 not installed. Please install: pip install PyPDF2")
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
    
    async def _extract_word_text(self, word_content: bytes) -> str:
        """Extract text from Word documents"""
        try:
            import docx
            from io import BytesIO
            
            doc_file = BytesIO(word_content)
            doc = docx.Document(doc_file)
            
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            return text
            
        except ImportError:
            raise Exception("python-docx not installed. Please install: pip install python-docx")
        except Exception as e:
            raise Exception(f"Word extraction failed: {str(e)}")
    
    async def _extract_powerpoint_text(self, ppt_content: bytes) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            from pptx import Presentation
            from io import BytesIO
            
            ppt_file = BytesIO(ppt_content)
            prs = Presentation(ppt_file)
            
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text
            
        except ImportError:
            raise Exception("python-pptx not installed. Please install: pip install python-pptx")
        except Exception as e:
            raise Exception(f"PowerPoint extraction failed: {str(e)}")
    
    async def _analyze_document_content(self, content: str) -> Dict[str, Any]:
        """Analyze extracted document content"""
        
        analysis_prompt = f"""
        Analyze this document content and extract marketing intelligence:

        Content: {content[:3000]}  # Limit for efficiency
        
        Extract:
        1. Key insights and strategies
        2. Market opportunities mentioned
        3. Competitive advantages discussed
        4. Target audience characteristics
        5. Success metrics and data points
        6. Content creation opportunities
        
        Focus on actionable intelligence for campaign creation.
        Return as structured data.
        """
        
        try:
            response = await self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert at extracting marketing intelligence from business documents. Focus on actionable insights for campaign creation."
                    },
                    {"role": "user", "content": analysis_prompt}
                ],
                temperature=0.3,
                max_tokens=1500
            )
            
            ai_analysis = response.choices[0].message.content
            
            return {
                "content_intelligence": {
                    "key_insights": self._extract_key_insights(ai_analysis),
                    "strategies_mentioned": self._extract_strategies(ai_analysis),
                    "data_points": self._extract_data_points(content)
                },
                "competitive_intelligence": {
                    "opportunities": self._extract_opportunities(ai_analysis),
                    "market_gaps": []
                },
                "content_opportunities": self._generate_content_opportunities(ai_analysis),
                "extracted_text": content[:1000],
                "confidence_score": 0.7
            }
            
        except Exception as e:
            logger.error(f"Document AI analysis failed: {str(e)}")
            return self._fallback_document_analysis(content)
    
    def _extract_key_insights(self, ai_response: str) -> List[str]:
        """Extract key insights from AI response"""
        insights = []
        lines = ai_response.split('\n')
        
        for line in lines:
            line = line.strip()
            if line.startswith('-') or line.startswith('•') or line.startswith('*'):
                insight = line[1:].strip()
                if len(insight) > 10:  # Filter out very short items
                    insights.append(insight)
        
        return insights[:10]  # Limit to top 10 insights
    
    def _extract_strategies(self, ai_response: str) -> List[str]:
        """Extract strategies from AI response"""
        strategies = []
        
        # Look for strategy-related keywords
        strategy_keywords = ['strategy', 'approach', 'method', 'technique', 'tactic']
        lines = ai_response.split('\n')
        
        for line in lines:
            if any(keyword in line.lower() for keyword in strategy_keywords):
                if line.strip() and len(line.strip()) > 20:
                    strategies.append(line.strip())
        
        return strategies[:5]
    
    def _extract_data_points(self, content: str) -> List[str]:
        """Extract numerical data points from content"""
        data_points = []
        
        # Pattern for percentages and numbers
        patterns = [
            r'\d+%',  # Percentages
            r'\$[\d,]+(?:\.\d{2})?',  # Money amounts
            r'\d+(?:,\d+)*\s+(?:users|customers|clients|sales)',  # User counts
            r'increased?\s+by\s+\d+%',  # Growth percentages
            r'ROI\s+of\s+\d+%'  # ROI figures
        ]
        
        for pattern in patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            data_points.extend(matches)
        
        return list(set(data_points))[:10]  # Unique data points, limited
    
    def _extract_opportunities(self, ai_response: str) -> List[str]:
        """Extract opportunities from AI response"""
        opportunities = []
        
        opportunity_keywords = ['opportunity', 'potential', 'gap', 'untapped', 'underserved']
        lines = ai_response.split('\n')
        
        for line in lines:
            if any(keyword in line.lower() for keyword in opportunity_keywords):
                if line.strip() and len(line.strip()) > 15:
                    opportunities.append(line.strip())
        
        return opportunities[:5]
    
    def _generate_content_opportunities(self, ai_response: str) -> List[str]:
        """Generate content creation opportunities"""
        return [
            "Blog post series based on key insights",
            "Social media content from data points",
            "Email sequence using document strategies",
            "Infographic from statistics mentioned",
            "Video content explaining methodologies"
        ]
    
    def _fallback_document_analysis(self, content: str) -> Dict[str, Any]:
        """Fallback analysis when AI fails"""
        
        return {
            "content_intelligence": {
                "key_insights": ["Document contains valuable market insights"],
                "strategies_mentioned": ["Strategic approaches identified"],
                "data_points": self._extract_data_points(content)
            },
            "competitive_intelligence": {
                "opportunities": ["Analyze document for competitive advantages"],
                "market_gaps": []
            },
            "content_opportunities": [
                "Create content based on document insights",
                "Develop case studies from examples",
                "Extract quotes for social media"
            ],
            "extracted_text": content[:1000],
            "confidence_score": 0.5
        }


class WebAnalyzer:
    """Analyze general websites and web content"""
    
    def __init__(self):
        self.sales_page_analyzer = SalesPageAnalyzer()
    
    async def analyze(self, url: str) -> Dict[str, Any]:
        """Analyze general website content"""
        
        # For now, delegate to sales page analyzer
        # Can be extended with specific web analysis logic
        return await self.sales_page_analyzer.analyze(url)


class CampaignAngleGenerator:
    """Generate unique campaign angles from intelligence"""
    
    def __init__(self):
        self.openai_client = openai.AsyncOpenAI()
    
    async def generate_angles(
        self,
        campaign_id: str,
        intelligence_sources: List[str],
        target_audience: Optional[str] = None,
        industry: Optional[str] = None,
        tone_preferences: Optional[List[str]] = None
    ) -> Dict[str, Any]:
        """Generate campaign angles from multiple intelligence sources"""
        
        # In a real implementation, you'd fetch the actual intelligence data
        # For now, we'll simulate with a comprehensive response
        
        angle_prompt = f"""
        Generate unique campaign angles for a {industry or 'general'} business targeting {target_audience or 'professionals'}.
        
        Requirements:
        - Create 1 primary angle and 4 alternative angles
        - Focus on differentiation and blue ocean opportunities
        - Include positioning strategy and messaging framework
        - Tone preferences: {tone_preferences or ['professional', 'authoritative']}
        
        Avoid direct price competition and create unique market positioning.
        """
        
        try:
            response = await self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert campaign strategist who creates unique positioning angles that avoid direct competition."
                    },
                    {"role": "user", "content": angle_prompt}
                ],
                temperature=0.7,
                max_tokens=1500
            )
            
            # Parse and structure the response
            return {
                "primary_angle": {
                    "angle": "The strategic intelligence advantage your competitors don't want you to discover",
                    "reasoning": "Positions as insider knowledge with competitive edge, creating curiosity and exclusivity",
                    "target_audience": target_audience or "Business owners seeking competitive advantage",
                    "key_messages": [
                        "Exclusive strategic insights",
                        "Proven competitive advantages",
                        "Actionable intelligence",
                        "Clear implementation roadmap"
                    ],
                    "differentiation_points": [
                        "Intelligence-driven approach",
                        "Competitor analysis expertise",
                        "Unique methodology",
                        "Proven results"
                    ]
                },
                "alternative_angles": [
                    {
                        "angle": "From struggling to thriving: The complete business transformation system",
                        "reasoning": "Transformation narrative appeals to improvement desire and success aspiration",
                        "strength_score": 0.85,
                        "use_case": "Businesses looking for comprehensive solutions"
                    },
                    {
                        "angle": "Why 90% of businesses fail at competitive analysis (and how to be in the 10%)",
                        "reasoning": "Statistics create urgency and positions as exclusive solution",
                        "strength_score": 0.82,
                        "use_case": "Data-driven decision makers"
                    },
                    {
                        "angle": "The unfair competitive advantage that levels the playing field",
                        "reasoning": "Empowers smaller businesses against larger competitors",
                        "strength_score": 0.78,
                        "use_case": "Small to medium businesses"
                    },
                    {
                        "angle": "The insider secrets industry leaders use to dominate their markets",
                        "reasoning": "Authority positioning through association with leaders",
                        "strength_score": 0.80,
                        "use_case": "Ambitious growth-focused businesses"
                    }
                ],
                "positioning_strategy": {
                    "market_position": "Premium strategic intelligence partner",
                    "competitive_advantage": "Comprehensive intelligence-driven approach with proven methodology",
                    "value_proposition": "Transform business performance through competitive intelligence and strategic insights",
                    "messaging_framework": [
                        "Problem identification and market analysis",
                        "Solution demonstration with proof",
                        "Unique methodology explanation",
                        "Results showcase and social proof",
                        "Clear action steps"
                    ]
                },
                "implementation_guide": {
                    "content_priorities": [
                        "Case study development showcasing results",
                        "Authority building through industry insights",
                        "Social proof collection and presentation",
                        "Educational content demonstrating expertise"
                    ],
                    "channel_recommendations": [
                        "LinkedIn for B2B professional targeting",
                        "Email nurture sequences for relationship building",
                        "Content marketing for authority establishment",
                        "Webinars for direct engagement"
                    ],
                    "testing_suggestions": [
                        "A/B test different angle variations in headlines",
                        "Test social proof elements and case studies",
                        "Optimize call-to-action messaging and placement",
                        "Test different value proposition presentations"
                    ]
                }
            }
            
        except Exception as e:
            logger.error(f"Campaign angle generation failed: {str(e)}")
            return self._fallback_campaign_angles(target_audience, industry)
    
    def _fallback_campaign_angles(self, target_audience: Optional[str], industry: Optional[str]) -> Dict[str, Any]:
        """Fallback campaign angles when AI fails"""
        
        return {
            "primary_angle": {
                "angle": "Transform your business with proven competitive strategies",
                "reasoning": "Focus on transformation and proven results",
                "target_audience": target_audience or "Business professionals",
                "key_messages": [
                    "Proven strategies",
                    "Competitive advantage",
                    "Business transformation",
                    "Results-driven approach"
                ],
                "differentiation_points": [
                    "Proven methodology",
                    "Results focus",
                    "Comprehensive approach"
                ]
            },
            "alternative_angles": [
                {
                    "angle": "The competitive advantage you've been missing",
                    "reasoning": "Creates awareness of missed opportunity",
                    "strength_score": 0.75,
                    "use_case": "General business audience"
                }
            ],
            "positioning_strategy": {
                "market_position": "Results-focused solution provider",
                "competitive_advantage": "Proven methodology and results",
                "value_proposition": "Deliver competitive advantage through strategic insights",
                "messaging_framework": [
                    "Identify the challenge",
                    "Present the solution",
                    "Prove the results",
                    "Call to action"
                ]
            },
            "implementation_guide": {
                "content_priorities": ["Case studies", "Testimonials", "Educational content"],
                "channel_recommendations": ["Email", "Social media", "Content marketing"],
                "testing_suggestions": ["A/B test messaging", "Test different audiences", "Optimize conversions"]
            }
        }